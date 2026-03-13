"""OneDrive/SharePoint connector via mounted filesystem.

Reads files directly from the macOS-mounted OneDrive/SharePoint folders —
no Graph API or authentication needed.
"""

import logging
import os
from datetime import datetime, timedelta, timezone
from pathlib import Path

from second_brain.config import get_config
from second_brain.connectors.base import BaseConnector
from second_brain.models import Item, ItemMeta

logger = logging.getLogger("second-brain.connectors.onedrive")

# File extensions we can extract text from
TEXT_EXTENSIONS = {".txt", ".md", ".csv", ".json", ".yaml", ".yml", ".xml", ".html"}
OFFICE_EXTENSIONS = {".docx", ".pptx", ".xlsx", ".pdf"}
SUPPORTED_EXTENSIONS = TEXT_EXTENSIONS | OFFICE_EXTENSIONS

# Default mount path on macOS
DEFAULT_ONEDRIVE_PATH = os.path.expanduser("~/Library/CloudStorage/OneDrive-Microsoft")


class OneDriveConnector(BaseConnector):
    """Connector that ingests documents from mounted OneDrive/SharePoint into the work store."""

    def __init__(self):
        super().__init__(name="OneDrive", target_store="work")
        self._config = get_config()
        self._base_path = Path(DEFAULT_ONEDRIVE_PATH)

    async def authenticate(self) -> bool:
        """Check that the mount point exists. No actual auth needed."""
        if self._base_path.is_dir():
            logger.info(f"OneDrive mounted at: {self._base_path}")
            return True
        logger.error(f"OneDrive not mounted at {self._base_path}")
        return False

    async def fetch_new_items(self, days_back: int = 7, subfolder: str | None = None) -> list[Item]:
        """Fetch recently modified documents from the mounted OneDrive/SharePoint.

        Args:
            days_back: How many days back to look for modified files.
            subfolder: Optional subfolder to scan (e.g. "MSHealthIL - Documents").
                       If None, scans the whole OneDrive.
        """
        scan_path = self._base_path / subfolder if subfolder else self._base_path
        if not scan_path.is_dir():
            logger.error(f"Path not found: {scan_path}")
            return []

        since = datetime.now(timezone.utc) - timedelta(days=days_back)
        since_ts = since.timestamp()
        items: list[Item] = []

        for root, _dirs, files in os.walk(scan_path):
            # Limit depth to 4 levels
            depth = Path(root).relative_to(scan_path).parts
            if len(depth) > 4:
                continue

            for filename in files:
                if filename.startswith(".") or filename.startswith("~$"):
                    continue

                filepath = Path(root, filename)
                ext = filepath.suffix.lower()
                if ext not in SUPPORTED_EXTENSIONS:
                    continue

                try:
                    stat = filepath.stat()
                except OSError:
                    continue

                if stat.st_mtime < since_ts:
                    continue

                modified = datetime.fromtimestamp(stat.st_mtime, tz=timezone.utc)
                rel_path = filepath.relative_to(self._base_path)

                content = self._read_file_content(filepath, ext)

                source_label = "sharepoint" if subfolder else "onedrive"
                item = Item(
                    meta=ItemMeta(
                        title=f"{filename}",
                        source=source_label,
                        store="work",
                        category="Inbox",
                        tags=[source_label, ext.lstrip(".")],
                        created=datetime.fromtimestamp(stat.st_birthtime, tz=timezone.utc)
                        if hasattr(stat, "st_birthtime")
                        else modified,
                        updated=modified,
                        summary=content[:200] if content else f"[{ext.upper()} file, {stat.st_size} bytes]",
                    ),
                    content=(
                        f"**File:** {rel_path}\n"
                        f"**Modified:** {modified.strftime('%Y-%m-%d %H:%M')}\n"
                        f"**Size:** {stat.st_size:,} bytes\n\n---\n\n{content}"
                    ),
                )
                items.append(item)

        logger.info(f"Found {len(items)} recently modified files in {scan_path.name}")
        return items

    def _read_file_content(self, filepath: Path, ext: str) -> str:
        """Read text content from a file. Returns empty string on failure."""
        try:
            if ext in TEXT_EXTENSIONS:
                return filepath.read_text(encoding="utf-8", errors="replace")[:10000]
            if ext == ".docx":
                return self._extract_docx(filepath)
            if ext == ".pptx":
                return self._extract_pptx(filepath)
            if ext == ".xlsx":
                return self._extract_xlsx(filepath)
            if ext == ".pdf":
                return self._extract_pdf(filepath)
            return ""
        except Exception as e:
            logger.debug(f"Content read failed for {filepath.name}: {e}")
            return ""

    @staticmethod
    def _extract_docx(filepath: Path) -> str:
        from docx import Document
        doc = Document(str(filepath))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paragraphs)[:10000]

    @staticmethod
    def _extract_pptx(filepath: Path) -> str:
        from pptx import Presentation
        prs = Presentation(str(filepath))
        parts: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            texts.append(text)
            if texts:
                parts.append(f"--- Slide {i} ---\n" + "\n".join(texts))
        return "\n\n".join(parts)[:10000]

    @staticmethod
    def _extract_xlsx(filepath: Path) -> str:
        from openpyxl import load_workbook
        wb = load_workbook(str(filepath), read_only=True, data_only=True)
        parts: list[str] = []
        for sheet_name in wb.sheetnames[:5]:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(max_row=50, values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    rows.append(" | ".join(cells))
            if rows:
                parts.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(parts)[:10000]

    @staticmethod
    def _extract_pdf(filepath: Path) -> str:
        from PyPDF2 import PdfReader
        reader = PdfReader(str(filepath))
        pages = []
        for i, page in enumerate(reader.pages[:20]):
            text = page.extract_text()
            if text and text.strip():
                pages.append(text.strip())
        return "\n\n".join(pages)[:10000]

    def list_folders(self) -> list[str]:
        """List top-level folders in the mounted OneDrive (for discovery)."""
        if not self._base_path.is_dir():
            return []
        return sorted(
            entry.name
            for entry in self._base_path.iterdir()
            if entry.is_dir() and not entry.name.startswith(".")
        )

    async def get_status(self) -> dict:
        """Get OneDrive connector status."""
        mounted = self._base_path.is_dir()
        folders = self.list_folders() if mounted else []
        return {
            "connector": "OneDrive/SharePoint (filesystem)",
            "mounted": mounted,
            "mount_path": str(self._base_path),
            "top_level_folders": folders,
            "target_store": self.target_store,
            "supported_formats": sorted(SUPPORTED_EXTENSIONS),
        }
